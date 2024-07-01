import base64
import io
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Generator, Literal, Union

from fastapi import FastAPI, UploadFile
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from requests import get


def emu_to_pixels(emu: int) -> float:
    return emu / 9525


def get_font_color(font) -> str:
    try:
        if font.color.rgb:
            return f"#{font.color.rgb:06x}"
    except AttributeError:
        pass
    return "#000000"  # Default color black if no color is set


def get_background_color(text_frame) -> str:
    try:
        if (
            text_frame.fill
            and text_frame.fill.fore_color
            and text_frame.fill.fore_color.rgb
        ):
            return f"#{text_frame.fill.fore_color.rgb:06x}"
    except AttributeError:
        pass
    return "transparent"  # Default background color


def get_slide_background_color(slide) -> str:
    try:
        fill = slide.background.fill
        if fill.type is not None:
            return f"#{fill.fore_color.rgb:06x}" if fill.fore_color.rgb else "#ccc"
    except (AttributeError, TypeError):
        pass
    return "#ccc"  # Default color if no valid background color is found


def stream_pptx(path: str) -> Generator[str, None, None]:
    prs = Presentation(path)
    yield "<!DOCTYPE html>"
    yield "<html>"
    yield """<head>
    <script src='https://cdn.tailwindcss.com'></script>
    <style>
    body { display: flex; flex-direction: column; justify-content: center; align-items: center; height: 100%; width: 100%; margin: 0; padding: 0; }
    section { display: flex; flex-direction: column; justify-content: center; align-items: center; height: 100vh; width: 100vw; position: relative; }
    img { width: 100%; height: auto; }
    p { position: absolute; word-wrap: break-word; white-space: normal; }
    footer { position: absolute; bottom: 10px; font-size: 12px; color: #666; }
    </style>
    </head>"""
    yield "<body>"
    for i, slide in enumerate(prs.slides, start=1):
        bg_color = get_slide_background_color(slide)
        yield f"<section id='slide-{i}' style='background-color: {bg_color};'>"
        for shape in slide.shapes:
            left = emu_to_pixels(shape.left)
            top = emu_to_pixels(shape.top)
            width = emu_to_pixels(shape.width)
            height = emu_to_pixels(shape.height)
            if shape.has_text_frame:
                text_frame = shape.text_frame
                font = text_frame.paragraphs[0].font
                font_size = font.size.pt if font.size else 12
                font_color = get_font_color(font)
                background_color = get_background_color(text_frame)
                bold = "font-bold" if font.bold else ""
                italic = "italic" if font.italic else ""
                underline = "underline" if font.underline else ""
                alignment = text_frame.paragraphs[0].alignment
                text_align = "left"
                if alignment == PP_ALIGN.CENTER:
                    text_align = "center"
                elif alignment == PP_ALIGN.RIGHT:
                    text_align = "right"
                yield f"<p style='left:{left}px; top:{top}px; width:{width}px; height:{height}px; font-size:{font_size}px; color:{font_color}; background-color:{background_color}; text-align:{text_align};' class='{bold} {italic} {underline} z-999'>{text_frame.text}</p>"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_base64 = base64.b64encode(image_bytes).decode()
                yield f"<img src='data:image/png;base64,{image_base64}' style='position:absolute; left:{left}px; top:{top}px; height:{height}px; width:{width}px;' />"
        yield f"<footer>Slide {i} of {len(prs.slides)}</footer>"
        yield "</section>"
    yield "</body>"
    yield "</html>"


@dataclass
class StreamPPTX:
    path: Union[str, Path, UploadFile, io.BytesIO, bytes]
    source_type: Literal["file", "url", "upload", "io", "bytes"]

    def handler(
        self,
        source: Union[str, Path, UploadFile, io.BytesIO, bytes],
        callback: Callable[[str], Generator[str, None, None]],
    ) -> StreamingResponse:
        with tempfile.NamedTemporaryFile(delete=False) as f:
            if self.source_type == "file" and isinstance(source, Path):
                f.write(source.read_bytes())
            elif self.source_type == "file" and isinstance(source, str):
                f.write(Path(source).read_bytes())
            elif self.source_type == "upload":
                assert isinstance(source, UploadFile)
                f.write(source.file.read())
            elif self.source_type == "io":
                assert isinstance(source, io.BytesIO)
                f.write(source.getvalue())
            elif self.source_type == "bytes":
                assert isinstance(source, bytes)
                f.write(source)
            elif self.source_type == "url":
                assert isinstance(source, str)
                response = get(source)
                f.write(response.content)
            else:
                raise ValueError(f"Invalid source type: {self.source_type}")
            self.path = f.name
        return StreamingResponse(callback(self.path), media_type="text/html")

    def run(self) -> StreamingResponse:
        return self.handler(self.path, stream_pptx)