from dataclasses import dataclass
from pptx import Presentation
from ._base import RawDoc


@dataclass
class PptxLoader(RawDoc):
    def extract_text(self):
        prs = Presentation(self.file_path)
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.has_text_frame:  # type: ignore
                    text_frame = shape.text_frame  # type: ignore
                    for paragraph in text_frame.paragraphs:  # type: ignore
                        if paragraph.text:  # type: ignore
                            yield paragraph.text
                        else:
                            continue

    def extract_image(self):
        prs = Presentation(self.file_path)
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.shape_type == 13:  # type: ignore
                    image = shape.image  # type: ignore
                    yield image.blob  # type: ignore
                else:
                    continue
